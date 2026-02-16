#!/usr/bin/env python3
"""
Blue Frontier - BLUE × PLC API Integration Presentation Generator
Generates a branded PowerPoint presentation highlighting the importance
of integrating the Distech PLC API with BLUE.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand Colors ─────────────────────────────────────────────────────────────
CYAN_BLUE    = RGBColor(0x00, 0xBF, 0xFF)
DEEP_BLUE    = RGBColor(0x1E, 0x3A, 0x8A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
MEDIUM_GRAY  = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
CRITICAL_RED = RGBColor(0xDC, 0x26, 0x26)
WARNING_YLW  = RGBColor(0xF5, 0x9E, 0x0B)
SUCCESS_GRN  = RGBColor(0x10, 0xB9, 0x81)
INFO_BLUE    = RGBColor(0x3B, 0x82, 0xF6)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Font names
HEADING_FONT = "Montserrat"
BODY_FONT    = "Avenir"


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape with optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_name=BODY_FONT,
                 font_size=14, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name=BODY_FONT,
                   font_size=14, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                   line_spacing=1.2, bullet=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        prefix = "\u2022  " if bullet else ""

        # Support bold prefix
        if isinstance(line, tuple):
            # (bold_part, normal_part)
            run1 = p.add_run()
            run1.text = prefix + line[0]
            run1.font.name = font_name
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = line[1]
            run2.font.name = font_name
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
        else:
            p.text = prefix + line
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold

        p.alignment = alignment
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_icon_card(slide, left, top, width, height, icon_text, title, body_lines,
                  accent_color=CYAN_BLUE):
    """Add a card with an icon circle, title, and bullet points."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False

    # Icon circle
    circle_size = Inches(0.7)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.3),
        top + Inches(0.3),
        circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()

    # Icon text
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(1.15), width - Inches(0.6),
                 Inches(0.4), title, font_name=HEADING_FONT, font_size=16,
                 color=DEEP_BLUE, bold=True)

    # Body bullets
    add_multi_text(slide, left + Inches(0.3), top + Inches(1.55), width - Inches(0.6),
                   height - Inches(1.8), body_lines, font_size=11, color=DARK_GRAY,
                   bullet=True, line_spacing=1.3)


def add_bottom_bar(slide):
    """Add a thin accent bar at the bottom of the slide."""
    add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08),
              SLIDE_WIDTH, Inches(0.08), CYAN_BLUE)


def add_slide_number(slide, num, total):
    """Add slide number in bottom-right corner."""
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.55),
                 Inches(1.2), Inches(0.3), f"{num} / {total}",
                 font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def add_section_header_bar(slide, text):
    """Add a colored header bar across the top with section title."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DEEP_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 text, font_name=HEADING_FONT, font_size=32, color=WHITE, bold=True)


# ═════════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    TOTAL_SLIDES = 16
    slide_num = 0
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DEEP_BLUE)

    # Accent stripe at top
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), CYAN_BLUE)

    # Main title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "BLUE × PLC API Integration",
                 font_name=HEADING_FONT, font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
                 "Giving BLUE Eyes on the Fleet",
                 font_name=HEADING_FONT, font_size=28, color=CYAN_BLUE, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Divider line
    add_shape(slide, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.04), CYAN_BLUE)

    # Company info
    add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.5),
                 "Blue Frontier  |  Distech ECY-600 PLC Integration Proposal",
                 font_name=BODY_FONT, font_size=16, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.5),
                 "February 2026",
                 font_name=BODY_FONT, font_size=14, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Try adding logo
    logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                             "..", "skills", "BF-Documentation-Skill", "assets", "BF-Logo-Dark.jpg")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.2), Inches(5.8), Inches(2.9))

    # Bottom accent bar
    add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.15), SLIDE_WIDTH, Inches(0.15), CYAN_BLUE)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2: The Opportunity
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "The Opportunity")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                 "BLUE is a brilliant consultant that needs to be fed data manually.",
                 font_name=BODY_FONT, font_size=22, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.8),
                 "With PLC API access, BLUE becomes an autonomous fleet management system that never sleeps.",
                 font_name=HEADING_FONT, font_size=24, color=DEEP_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Two-column comparison
    # LEFT: Without API
    left_card = add_rounded_rect(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(3.2),
                                  RGBColor(0xFF, 0xF0, 0xF0))
    add_text_box(slide, Inches(1.1), Inches(3.75), Inches(5), Inches(0.5),
                 "Without API Access", font_name=HEADING_FONT, font_size=18,
                 color=CRITICAL_RED, bold=True)
    add_multi_text(slide, Inches(1.1), Inches(4.3), Inches(5), Inches(2.3),
                   [
                       "Requires manual data uploads to analyze issues",
                       "Cannot verify what the controller is actually doing",
                       "Reactive — responds only after problems are reported",
                       "No real-time visibility into deployed fleet",
                       "Troubleshooting relies on technician-provided data",
                   ], font_size=13, color=DARK_GRAY, bullet=True, line_spacing=1.4)

    # RIGHT: With API
    right_card = add_rounded_rect(slide, Inches(7), Inches(3.6), Inches(5.5), Inches(3.2),
                                   RGBColor(0xEF, 0xFF, 0xEF))
    add_text_box(slide, Inches(7.3), Inches(3.75), Inches(5), Inches(0.5),
                 "With API Access", font_name=HEADING_FONT, font_size=18,
                 color=SUCCESS_GRN, bold=True)
    add_multi_text(slide, Inches(7.3), Inches(4.3), Inches(5), Inches(2.3),
                   [
                       "Queries PLCs directly for live data",
                       "Verifies actual vs. expected controller behavior",
                       "Proactive — detects issues before they escalate",
                       "24/7 autonomous fleet monitoring",
                       "Guides techs with real-time, data-backed insight",
                   ], font_size=13, color=DARK_GRAY, bullet=True, line_spacing=1.4)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3: Current BLUE Capabilities
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "What BLUE Can Do Today")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(1), Inches(1.4), Inches(11), Inches(0.6),
                 "EC-GFX Mastery Skill — Expert-Level PLC Understanding",
                 font_name=HEADING_FONT, font_size=20, color=DEEP_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    cards = [
        ("LB", "Logic Blocks", [
            "Interpret and troubleshoot",
            "all EC-GFX logic blocks",
            "including custom sequences"
        ]),
        ("PID", "Control Loops", [
            "Analyze PID loops,",
            "control sequences,",
            "and state machines"
        ]),
        ("AL", "Alarm Analysis", [
            "Diagnose alarm conditions,",
            "error states, and",
            "latched fault logic"
        ]),
        ("BP", "Best Practices", [
            "Programming standards,",
            "optimization guidance,",
            "and technician training"
        ]),
    ]

    for i, (icon, title, body) in enumerate(cards):
        x = Inches(0.6 + i * 3.15)
        add_icon_card(slide, x, Inches(2.2), Inches(2.85), Inches(3.2),
                      icon, title, body, accent_color=CYAN_BLUE)

    add_text_box(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.8),
                 "These skills are powerful — but they operate on static data.\n"
                 "BLUE can only analyze what is manually provided.",
                 font_name=BODY_FONT, font_size=15, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4: The Blind Spots
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "What BLUE Cannot See — Yet")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    blind_spots = [
        ("Real-time access to actual PLC data", "No live sensor readings, valve positions, or actuator states"),
        ("Live point values and setpoints", "Cannot see what the controller is targeting right now"),
        ("Historical trending from controllers", "Must rely on Grafana exports instead of direct access"),
        ("Verification of controller behavior", "Cannot compare actual vs. expected execution"),
        ("Direct visibility into logic execution", "Blind to what the PLC is actually doing in real time"),
    ]

    for i, (title, desc) in enumerate(blind_spots):
        y = Inches(1.5 + i * 1.1)
        # Red X indicator
        x_mark = add_rounded_rect(slide, Inches(1), y, Inches(0.5), Inches(0.5), CRITICAL_RED)
        tf = x_mark.text_frame
        p = tf.paragraphs[0]
        p.text = "X"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(1.8), y - Inches(0.05), Inches(10), Inches(0.35),
                     title, font_name=HEADING_FONT, font_size=17, color=DEEP_BLUE, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.3), Inches(10), Inches(0.35),
                     desc, font_name=BODY_FONT, font_size=13, color=MEDIUM_GRAY)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5: What the API Unlocks (Overview)
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DEEP_BLUE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_BLUE)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                 "Six Capabilities the API Unlocks",
                 font_name=HEADING_FONT, font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    capabilities = [
        ("01", "Real-Time Diagnostics", "Query PLCs directly for live status and sensor data"),
        ("02", "Automated Monitoring", "24/7 fleet polling with automatic anomaly detection"),
        ("03", "Remote Troubleshooting", "Guide technicians with live PLC state data"),
        ("04", "Predictive Intelligence", "Detect patterns that precede failures"),
        ("05", "Control Optimization", "Analyze and improve actual control performance"),
        ("06", "Verification & Validation", "Confirm repairs and setpoint changes objectively"),
    ]

    for i, (num, title, desc) in enumerate(capabilities):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(2.0 + row * 2.6)

        card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.2),
                                 RGBColor(0x16, 0x42, 0x96))  # Slightly lighter deep blue

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2),
                                        Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CYAN_BLUE
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.5),
                     title, font_name=HEADING_FONT, font_size=18, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.4), Inches(0.7),
                     desc, font_name=BODY_FONT, font_size=13, color=LIGHT_GRAY)

    add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), CYAN_BLUE)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6: Real-Time Diagnostics (Deep Dive)
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "01  Real-Time Diagnostics")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
                 'Instead of saying "check if the regenerator is running," BLUE could tell you exactly what is happening.',
                 font_name=BODY_FONT, font_size=16, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    examples = [
        ("Query PLC Directly",
         '"Regenerator status is OFF, last cycle completed 47 minutes ago"',
         "Instant, definitive answers instead of guesswork"),
        ("Live Sensor Readings",
         "See exactly what every sensor is reading right now",
         "Temperature, pressure, flow, concentration — all live"),
        ("Instant Failure Detection",
         "Identify stuck values, out-of-range readings, sensor failures",
         "Catch problems the moment they occur, not hours later"),
    ]

    for i, (title, example, desc) in enumerate(examples):
        y = Inches(2.3 + i * 1.6)
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.35), LIGHT_GRAY)

        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(11), Inches(0.35),
                     title, font_name=HEADING_FONT, font_size=16, color=DEEP_BLUE, bold=True)

        # Example in cyan box
        example_box = add_rounded_rect(slide, Inches(1.2), y + Inches(0.5), Inches(10.5),
                                        Inches(0.4), RGBColor(0xE0, 0xF7, 0xFF))
        add_text_box(slide, Inches(1.4), y + Inches(0.5), Inches(10), Inches(0.4),
                     example, font_name=BODY_FONT, font_size=12, color=DEEP_BLUE, bold=True)

        add_text_box(slide, Inches(1.2), y + Inches(0.95), Inches(11), Inches(0.3),
                     desc, font_name=BODY_FONT, font_size=12, color=MEDIUM_GRAY)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Automated Monitoring
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "02  Automated Monitoring")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
                 "Autonomous agents that watch the fleet around the clock",
                 font_name=BODY_FONT, font_size=18, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    monitoring_items = [
        ("15", "min", "Polling Interval",
         "Every deployed unit queried automatically, continuously building a real-time picture of fleet health"),
        ("24/7", "", "Coverage",
         "No weekends, no holidays, no fatigue — BLUE monitors every unit without interruption"),
        ("< 1", "min", "Alert Response",
         "Detect control logic failures, setpoint drift, and overrides before they impact performance"),
    ]

    for i, (big_num, unit, title, desc) in enumerate(monitoring_items):
        x = Inches(0.6 + i * 4.2)
        y = Inches(2.3)

        card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.2), LIGHT_GRAY)

        # Big metric
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(1.0),
                     big_num, font_name=HEADING_FONT, font_size=54, color=CYAN_BLUE,
                     bold=True, alignment=PP_ALIGN.CENTER)
        if unit:
            add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.4),
                         unit, font_name=BODY_FONT, font_size=16, color=MEDIUM_GRAY,
                         alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.2), Inches(0.4),
                     title, font_name=HEADING_FONT, font_size=18, color=DEEP_BLUE,
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Divider
        add_shape(slide, x + Inches(1.2), y + Inches(2.2), Inches(1.4), Inches(0.03), CYAN_BLUE)

        add_text_box(slide, x + Inches(0.3), y + Inches(2.5), Inches(3.2), Inches(1.5),
                     desc, font_name=BODY_FONT, font_size=12, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8: Remote Troubleshooting
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "03  Remote Troubleshooting")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
                 "When a technician calls with an issue, BLUE provides real-time intelligence — not guesswork",
                 font_name=BODY_FONT, font_size=16, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    steps = [
        ("1", "Pull Up PLC State", "Instantly retrieve the exact\ncontroller state in real time"),
        ("2", "Analyze Controller", "See what the controller thinks\nis happening vs. reality"),
        ("3", "Identify Root Cause", "Pinpoint logic errors or sensor\nproblems immediately"),
        ("4", "Guide the Tech", "Provide step-by-step guidance\nbacked by live data"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        y = Inches(2.5)

        # Step card
        card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(3.5), LIGHT_GRAY)

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.25),
                                         Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = DEEP_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.5), Inches(0.5),
                     title, font_name=HEADING_FONT, font_size=16, color=DEEP_BLUE,
                     bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.5), Inches(1.2),
                     desc, font_name=BODY_FONT, font_size=13, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

        # Arrow between cards
        if i < 3:
            arrow_x = x + Inches(3.0)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(1.5),
                                            Inches(0.25), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CYAN_BLUE
            arrow.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
                 "Result: Faster resolution, fewer truck rolls, higher first-time fix rates",
                 font_name=HEADING_FONT, font_size=16, color=SUCCESS_GRN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Predictive Intelligence
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "04  Predictive Intelligence")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
                 "Move from reactive service to predictive fleet management",
                 font_name=BODY_FONT, font_size=18, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    pred_items = [
        ("Control Stability", [
            "Track excessive cycling and hunting",
            "Detect oscillating PID loops",
            "Identify control instability early",
        ]),
        ("Failure Prediction", [
            "Recognize patterns that precede failures",
            "Alert before equipment breaks down",
            "Reduce unplanned downtime",
        ]),
        ("PID Tuning", [
            "Detect poorly tuned loops",
            "Recommend tuning adjustments",
            "Based on actual real-world behavior",
        ]),
        ("Comm Health", [
            "Monitor controller communication",
            "Detect network degradation",
            "Track response times and errors",
        ]),
    ]

    for i, (title, bullets) in enumerate(pred_items):
        x = Inches(0.5 + i * 3.2)
        y = Inches(2.3)

        card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(3.8), LIGHT_GRAY)

        # Accent bar at top of card
        add_shape(slide, x, y, Inches(2.9), Inches(0.06), CYAN_BLUE)

        add_text_box(slide, x + Inches(0.25), y + Inches(0.3), Inches(2.4), Inches(0.4),
                     title, font_name=HEADING_FONT, font_size=16, color=DEEP_BLUE,
                     bold=True, alignment=PP_ALIGN.CENTER)

        add_multi_text(slide, x + Inches(0.25), y + Inches(0.9), Inches(2.4), Inches(2.5),
                       bullets, font_size=13, color=DARK_GRAY, bullet=True, line_spacing=1.5)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10: Control Optimization
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "05  Control Optimization")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    opt_items = [
        ("Performance Analysis",
         "Compare actual control performance against optimal sequences to identify gaps and inefficiencies"),
        ("PID Tuning Recommendations",
         "Suggest tuning adjustments based on real behavior patterns observed across the fleet"),
        ("Strategy Improvement",
         "Identify inefficient control strategies and recommend improvements backed by data"),
        ("Fleet-Wide Insights",
         "Leverage data from every deployed unit to optimize logic for the entire fleet simultaneously"),
    ]

    for i, (title, desc) in enumerate(opt_items):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.6 + row * 2.8)

        card = add_rounded_rect(slide, x, y, Inches(5.9), Inches(2.4), LIGHT_GRAY)

        # Numbered circle
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3),
                                        Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CYAN_BLUE
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, x + Inches(1.2), y + Inches(0.3), Inches(4.4), Inches(0.45),
                     title, font_name=HEADING_FONT, font_size=18, color=DEEP_BLUE, bold=True)
        add_text_box(slide, x + Inches(1.2), y + Inches(0.9), Inches(4.4), Inches(1.3),
                     desc, font_name=BODY_FONT, font_size=14, color=DARK_GRAY)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 11: Verification & Validation
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "06  Verification & Validation")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
                 "Objective, data-driven confirmation that work was done correctly",
                 font_name=BODY_FONT, font_size=18, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    vv_items = [
        ("Verify Repairs", "Check actual PLC response after service work to confirm the fix holds"),
        ("Confirm Setpoints", "Validate that setpoint changes took effect and are being maintained"),
        ("Validate Sequences", "Ensure control sequences execute correctly through full cycles"),
        ('Job Complete" Confirmation', "Provide objective, data-backed proof that the job is done right"),
    ]

    for i, (title, desc) in enumerate(vv_items):
        y = Inches(2.3 + i * 1.2)

        # Check mark circle
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.55), Inches(0.55))
        check.fill.solid()
        check.fill.fore_color.rgb = SUCCESS_GRN
        check.line.fill.background()
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "\u2713"
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(1.8), y - Inches(0.02), Inches(10), Inches(0.35),
                     title, font_name=HEADING_FONT, font_size=17, color=DEEP_BLUE, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(10), Inches(0.35),
                     desc, font_name=BODY_FONT, font_size=14, color=DARK_GRAY)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 12: Integration Architecture
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "Integration Architecture")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    # Architecture flow - vertical
    arch_layers = [
        ("Blue Frontier Fleet", "All deployed LD-DOAS units in the field", DEEP_BLUE, WHITE),
        ("Distech PLCs (API Endpoints)", "ECY-600 controllers exposing RESTful API", CYAN_BLUE, WHITE),
        ("BLUE Monitoring Agents", "Autonomous polling every 15 minutes", RGBColor(0x16, 0x42, 0x96), WHITE),
        ("Data Analysis & Pattern Recognition", "Anomaly detection, trending, fleet comparison", INFO_BLUE, WHITE),
    ]

    for i, (title, desc, bg_color, text_color) in enumerate(arch_layers):
        y = Inches(1.3 + i * 1.2)
        box = add_rounded_rect(slide, Inches(3.5), y, Inches(6.3), Inches(0.9), bg_color)
        add_text_box(slide, Inches(3.7), y + Inches(0.05), Inches(5.9), Inches(0.4),
                     title, font_name=HEADING_FONT, font_size=16, color=text_color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3.7), y + Inches(0.45), Inches(5.9), Inches(0.35),
                     desc, font_name=BODY_FONT, font_size=11, color=RGBColor(0xDD, 0xDD, 0xFF),
                     alignment=PP_ALIGN.CENTER)

        # Arrow between layers
        if i < len(arch_layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                            Inches(6.5), y + Inches(0.9),
                                            Inches(0.35), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CYAN_BLUE
            arrow.line.fill.background()

    # Automated Actions - branching out
    actions_y = Inches(5.3)
    add_text_box(slide, Inches(3.5), actions_y - Inches(0.1), Inches(6.3), Inches(0.4),
                 "Automated Actions", font_name=HEADING_FONT, font_size=14, color=DEEP_BLUE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    actions = [
        ("Alert Generation", SUCCESS_GRN),
        ("ClickUp Tickets", INFO_BLUE),
        ("Trend Analysis", WARNING_YLW),
        ("Performance Reports", CYAN_BLUE),
    ]

    for i, (action, color) in enumerate(actions):
        x = Inches(2.4 + i * 2.3)
        box = add_rounded_rect(slide, x, actions_y + Inches(0.35), Inches(2.0), Inches(0.65), color)
        add_text_box(slide, x + Inches(0.1), actions_y + Inches(0.4), Inches(1.8), Inches(0.55),
                     action, font_name=BODY_FONT, font_size=12, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 13: What We Need
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "What BLUE Needs to Get Started")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    needs = [
        ("API Documentation", [
            "Endpoint URLs and authentication methods",
            "Available data points and their formats",
            "Read/write capabilities",
            "Polling rate limits",
        ]),
        ("Sample Data", [
            "Example API responses",
            "Point naming conventions",
            "Units of measurement for each point",
        ]),
        ("Access Credentials", [
            "Authentication method (API keys, OAuth, etc.)",
            "Per-unit vs. shared credentials",
            "Security token management",
        ]),
        ("Network Architecture", [
            "Connection method (direct, VPN, gateway)",
            "Firewall / security considerations",
            "Latency expectations",
        ]),
    ]

    for i, (title, items) in enumerate(needs):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.5 + row * 2.8)

        card = add_rounded_rect(slide, x, y, Inches(5.9), Inches(2.5), LIGHT_GRAY)

        # Header bar in card
        add_shape(slide, x, y, Inches(5.9), Inches(0.55), DEEP_BLUE)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.08), Inches(5.3), Inches(0.4),
                     title, font_name=HEADING_FONT, font_size=16, color=WHITE, bold=True)

        add_multi_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(5.3), Inches(1.6),
                       items, font_size=13, color=DARK_GRAY, bullet=True, line_spacing=1.4)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 14: Implementation Roadmap
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_section_header_bar(slide, "Implementation Roadmap — Day One")
    add_bottom_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    phases = [
        ("Phase 1", "API Connection\n& Validation", "Morning",
         ["Authenticate and connect to first PLC",
          "Map all available data points",
          "Verify data quality and refresh rates",
          "Test read capabilities"],
         DEEP_BLUE),
        ("Phase 2", "Monitoring\nAgent Build", "Midday",
         ["Build autonomous polling agent",
          "Set up data storage and trending",
          "Create alert conditions",
          "Test on 1-2 pilot units"],
         CYAN_BLUE),
        ("Phase 3", "Integration\n& Deployment", "Afternoon",
         ["Connect to ClickUp for auto-ticketing",
          "Deploy across fleet",
          "Establish baseline metrics",
          "Begin 24/7 monitoring"],
         SUCCESS_GRN),
    ]

    for i, (phase, title, time, steps, color) in enumerate(phases):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.5)

        # Phase header
        header = add_rounded_rect(slide, x, y, Inches(3.8), Inches(1.4), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.3),
                     phase + "  |  " + time, font_name=BODY_FONT, font_size=11,
                     color=RGBColor(0xDD, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.4), Inches(0.9),
                     title, font_name=HEADING_FONT, font_size=20, color=WHITE,
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Steps card
        card = add_rounded_rect(slide, x, Inches(3.1), Inches(3.8), Inches(3.5), LIGHT_GRAY)
        add_multi_text(slide, x + Inches(0.3), Inches(3.3), Inches(3.2), Inches(3.0),
                       steps, font_size=13, color=DARK_GRAY, bullet=True, line_spacing=1.5)

        # Arrow between phases
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + Inches(3.9), Inches(2.0),
                                            Inches(0.25), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 15: Day One Goals
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DEEP_BLUE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_BLUE)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)

    add_text_box(slide, Inches(1), Inches(0.7), Inches(11), Inches(0.8),
                 "Day One Deliverables",
                 font_name=HEADING_FONT, font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
                 "By end of day, BLUE will have:",
                 font_name=BODY_FONT, font_size=18, color=CYAN_BLUE,
                 alignment=PP_ALIGN.CENTER)

    goals = [
        "Live connection to at least one deployed unit",
        "Real-time dashboard showing actual PLC data",
        "At least one automated alert working",
        "Proof of concept for autonomous fleet monitoring",
    ]

    for i, goal in enumerate(goals):
        y = Inches(2.4 + i * 1.15)
        card = add_rounded_rect(slide, Inches(2), y, Inches(9.3), Inches(0.9),
                                 RGBColor(0x16, 0x42, 0x96))

        # Check circle
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.3), y + Inches(0.15),
                                        Inches(0.6), Inches(0.6))
        check.fill.solid()
        check.fill.fore_color.rgb = CYAN_BLUE
        check.line.fill.background()
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(3.2), y + Inches(0.15), Inches(7.8), Inches(0.6),
                     goal, font_name=HEADING_FONT, font_size=20, color=WHITE, bold=False)

    add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), CYAN_BLUE)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 16: Closing
    # ─────────────────────────────────────────────────────────────────────────
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DEEP_BLUE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), CYAN_BLUE)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
                 "Let's Give BLUE Eyes on the Fleet",
                 font_name=HEADING_FONT, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.04), CYAN_BLUE)

    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.2),
                 "The difference is like having a technician who can look at documentation\n"
                 "versus a technician who can actually see what the equipment is doing.",
                 font_name=BODY_FONT, font_size=20, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
                 "Share the API documentation. Let's build this integration.",
                 font_name=HEADING_FONT, font_size=22, color=CYAN_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Logo
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.2), Inches(5.8), Inches(2.9))

    add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.15), SLIDE_WIDTH, Inches(0.15), CYAN_BLUE)

    # ─── SAVE ────────────────────────────────────────────────────────────────
    output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(output_dir, "BLUE_PLC_API_Integration_Proposal.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
